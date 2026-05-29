#!/usr/bin/env python3
"""Inject PowerPoint animations into a pptx based on a JSON spec.

This script does two things, in order of reliability:

1. ALWAYS mirrors the animation plan into each slide's speaker notes.
   This is the high-value, always-correct output — a human can read the
   notes in PowerPoint and verify/apply the animation order manually.

2. BEST-EFFORT injects OOXML animation timing nodes for the supported
   subset of animations (entrance: fade/appear/fly in/wipe/zoom, with
   click and after-previous triggers). When an animation falls outside
   the supported set, the script logs it and relies on the speaker-note
   mirror as the fallback path.

PowerPoint animation XML is brittle across versions. Treat the speaker
notes as the source of truth and the XML injection as a convenience.

Spec format (animations.json):

    {
      "slides": [
        {
          "slide_index": 0,
          "steps": [
            {
              "target": "T1",
              "animation": "fade_in",
              "trigger": "on_load",
              "duration_ms": 500
            },
            {
              "target": "G1",
              "animation": "fly_in_left",
              "trigger": "on_click",
              "duration_ms": 400
            },
            {
              "target": "A1",
              "animation": "wipe_right",
              "trigger": "after_prev",
              "delay_ms": 200,
              "duration_ms": 300
            }
          ]
        }
      ]
    }

Usage:

    python inject_animations.py --in deck.pptx --spec animations.json --out deck_animated.pptx
"""

from __future__ import annotations

import argparse
import json
import sys
from dataclasses import dataclass
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn


SUPPORTED_ANIMATIONS = {
    "appear",
    "fade_in",
    "fade_out",
    "fly_in_left",
    "fly_in_right",
    "fly_in_top",
    "fly_in_bottom",
    "wipe_left",
    "wipe_right",
    "wipe_up",
    "wipe_down",
    "zoom",
}

# Triggers that the OOXML stCond mapping handles reliably. `after_prev`
# and `with_prev` resolve their parent-relative semantics through PPT's
# sequence interpretation, which is fragile across versions when authored
# by hand — we route those steps to speaker-notes-only and let the human
# set them in PowerPoint's Animation Pane.
SUPPORTED_TRIGGERS = {"on_load", "on_click"}

# PowerPoint preset class / preset id mapping for entrance animations.
# Reference: ECMA-376 Part 1 §19.5 and Microsoft Open XML SDK examples.
PRESET_IDS = {
    "appear":         {"class": "entr", "presetID": "1",  "presetSubtype": "0"},
    "fade_in":        {"class": "entr", "presetID": "10", "presetSubtype": "0"},
    "fade_out":       {"class": "exit", "presetID": "10", "presetSubtype": "0"},
    "fly_in_left":    {"class": "entr", "presetID": "2",  "presetSubtype": "4"},
    "fly_in_right":   {"class": "entr", "presetID": "2",  "presetSubtype": "8"},
    "fly_in_top":     {"class": "entr", "presetID": "2",  "presetSubtype": "1"},
    "fly_in_bottom":  {"class": "entr", "presetID": "2",  "presetSubtype": "2"},
    "wipe_left":      {"class": "entr", "presetID": "12", "presetSubtype": "4"},
    "wipe_right":     {"class": "entr", "presetID": "12", "presetSubtype": "8"},
    "wipe_up":        {"class": "entr", "presetID": "12", "presetSubtype": "1"},
    "wipe_down":      {"class": "entr", "presetID": "12", "presetSubtype": "2"},
    "zoom":           {"class": "entr", "presetID": "23", "presetSubtype": "0"},
}


@dataclass
class StepResult:
    target: str
    animation: str
    injected_xml: bool
    note: str


def human_readable_step(step: dict, index: int) -> str:
    """Format a step for inclusion in speaker notes."""
    trigger_map = {
        "on_load": "On slide load",
        "on_click": "On click",
        "after_prev": "After previous",
        "with_prev": "With previous",
    }
    trigger = trigger_map.get(step["trigger"], step["trigger"])
    delay = step.get("delay_ms", 0)
    if delay and step["trigger"] == "after_prev":
        trigger = f"After previous +{delay / 1000:.1f}s"
    duration_s = step["duration_ms"] / 1000
    return (
        f"{index + 1}. {step['target']} — "
        f"{step['animation'].replace('_', ' ')} — "
        f"{trigger} ({duration_s:.1f}s)"
    )


def write_speaker_notes(slide, steps: list[dict]) -> None:
    """Mirror the animation plan into the slide's speaker notes.

    This is always done, regardless of whether XML injection succeeds.
    """
    lines = ["ANIMATION ORDER:"]
    lines.extend(human_readable_step(step, i) for i, step in enumerate(steps))
    notes_text = "\n".join(lines)

    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    existing = text_frame.text.strip()
    if existing and "ANIMATION ORDER:" not in existing:
        text_frame.text = f"{existing}\n\n{notes_text}"
    else:
        text_frame.text = notes_text


def find_shape_by_name(slide, name: str):
    """Locate a shape (or group) on the slide by its name attribute.

    Returns the shape object or None if not found.
    """
    for shape in slide.shapes:
        if shape.name == name:
            return shape
        # Check nested group members
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            for member in shape.shapes:
                if member.name == name:
                    return member
    return None


def build_timing_xml(steps: list[dict], shape_id_lookup: dict[str, int]) -> etree._Element:
    """Build the <p:timing> XML element for a slide's animation sequence.

    This generates a "main sequence" timing block — the standard
    click-through reveal pattern used for ~90% of PPT animations.

    Steps with unsupported animation types are silently skipped (their
    speaker-note mirror is the fallback).
    """
    nsmap = {
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    }
    p_ns = "{http://schemas.openxmlformats.org/presentationml/2006/main}"

    timing = etree.SubElement(etree.Element("root"), f"{p_ns}timing", nsmap=nsmap)
    timing.getparent().remove(timing)  # detach from temp root

    tn_lst = etree.SubElement(timing, f"{p_ns}tnLst")
    root_par = etree.SubElement(tn_lst, f"{p_ns}par")
    root_ctn = etree.SubElement(root_par, f"{p_ns}cTn", {
        "id": "1", "dur": "indefinite", "restart": "never", "nodeType": "tmRoot",
    })
    root_children = etree.SubElement(root_ctn, f"{p_ns}childTnLst")

    seq = etree.SubElement(root_children, f"{p_ns}seq", {
        "concurrent": "1", "nextAc": "seek",
    })
    seq_ctn = etree.SubElement(seq, f"{p_ns}cTn", {
        "id": "2", "dur": "indefinite", "nodeType": "mainSeq",
    })
    seq_children = etree.SubElement(seq_ctn, f"{p_ns}childTnLst")

    node_id = 3
    for step in steps:
        anim = step["animation"]
        if anim not in PRESET_IDS:
            continue
        target = step["target"]
        if target not in shape_id_lookup:
            continue
        sp_id = shape_id_lookup[target]
        preset = PRESET_IDS[anim]
        duration = step["duration_ms"]
        trigger = step["trigger"]
        delay = step.get("delay_ms", 0)

        # Build one click-effect entry (p:par > p:cTn > p:childTnLst > p:par > p:cTn ...)
        click_par = etree.SubElement(seq_children, f"{p_ns}par")
        click_ctn = etree.SubElement(click_par, f"{p_ns}cTn", {
            "id": str(node_id),
            "fill": "hold",
        })
        node_id += 1

        # Trigger condition. "on_load" fires when the slide opens (delay 0
        # from slide-begin); "on_click" waits for the user's next click.
        st_cond_lst = etree.SubElement(click_ctn, f"{p_ns}stCondLst")
        cond_evt = "onBegin" if trigger == "on_load" else "onClick"
        etree.SubElement(st_cond_lst, f"{p_ns}cond", {
            "evt": cond_evt, "delay": str(delay),
        })

        click_children = etree.SubElement(click_ctn, f"{p_ns}childTnLst")
        effect_par = etree.SubElement(click_children, f"{p_ns}par")
        effect_ctn = etree.SubElement(effect_par, f"{p_ns}cTn", {
            "id": str(node_id),
            "presetID": preset["presetID"],
            "presetClass": preset["class"],
            "presetSubtype": preset["presetSubtype"],
            "fill": "hold",
            "grpId": "0",
            "nodeType": "clickEffect" if trigger == "on_click" else "withEffect",
        })
        node_id += 1

        effect_children = etree.SubElement(effect_ctn, f"{p_ns}childTnLst")
        set_node = etree.SubElement(effect_children, f"{p_ns}set")
        c_behavior = etree.SubElement(set_node, f"{p_ns}cBhvr")
        # Duration belongs on the behavior cTn — PowerPoint reads animation
        # length from here, not from the outer effect container.
        c_tn = etree.SubElement(c_behavior, f"{p_ns}cTn", {
            "id": str(node_id),
            "dur": str(duration),
            "fill": "hold",
        })
        node_id += 1
        st_cond = etree.SubElement(c_tn, f"{p_ns}stCondLst")
        etree.SubElement(st_cond, f"{p_ns}cond", {"delay": "0"})
        tgt_el = etree.SubElement(c_behavior, f"{p_ns}tgtEl")
        etree.SubElement(tgt_el, f"{p_ns}spTgt", {"spid": str(sp_id)})
        attr_name_lst = etree.SubElement(c_behavior, f"{p_ns}attrNameLst")
        attr_name = etree.SubElement(attr_name_lst, f"{p_ns}attrName")
        attr_name.text = "style.visibility"
        to_node = etree.SubElement(set_node, f"{p_ns}to")
        etree.SubElement(to_node, f"{p_ns}strVal", {"val": "visible"})

    # Sequence-level conditions: which clicks advance the sequence
    prev_cond_lst = etree.SubElement(seq, f"{p_ns}prevCondLst")
    prev_cond = etree.SubElement(prev_cond_lst, f"{p_ns}cond", {
        "evt": "onPrev", "delay": "0",
    })
    prev_tgt = etree.SubElement(prev_cond, f"{p_ns}tgtEl")
    etree.SubElement(prev_tgt, f"{p_ns}sldTgt")

    next_cond_lst = etree.SubElement(seq, f"{p_ns}nextCondLst")
    next_cond = etree.SubElement(next_cond_lst, f"{p_ns}cond", {
        "evt": "onNext", "delay": "0",
    })
    next_tgt = etree.SubElement(next_cond, f"{p_ns}tgtEl")
    etree.SubElement(next_tgt, f"{p_ns}sldTgt")

    return timing


def collect_shape_ids(slide) -> dict[str, int]:
    """Map shape name → shape XML id for the slide.

    Animation XML references shapes by their numeric `id` attribute,
    not by name. We map name → id so the user's plan IDs (T1, G2, ...)
    can be resolved.
    """
    lookup: dict[str, int] = {}

    def walk(shape):
        sp_pr = shape._element
        nv_sp_pr = sp_pr.find(qn("p:nvSpPr"))
        if nv_sp_pr is None:
            nv_sp_pr = sp_pr.find(qn("p:nvGrpSpPr"))
        if nv_sp_pr is None:
            nv_sp_pr = sp_pr.find(qn("p:nvPicPr"))
        if nv_sp_pr is None:
            nv_sp_pr = sp_pr.find(qn("p:nvCxnSpPr"))
        if nv_sp_pr is None:
            return
        c_nv_pr = nv_sp_pr.find(qn("p:cNvPr"))
        if c_nv_pr is None:
            return
        name = c_nv_pr.get("name")
        shape_id = int(c_nv_pr.get("id", "0"))
        if name and shape_id:
            lookup[name] = shape_id

        if shape.shape_type == 6:  # group
            for member in shape.shapes:
                walk(member)

    for s in slide.shapes:
        walk(s)
    return lookup


def inject_for_slide(slide, steps: list[dict]) -> list[StepResult]:
    """Inject animation XML for one slide. Returns per-step results."""
    results: list[StepResult] = []

    # Always write speaker notes first — this is the reliable path.
    write_speaker_notes(slide, steps)

    shape_id_lookup = collect_shape_ids(slide)

    injectable_steps = []
    for step in steps:
        anim = step["animation"]
        target = step["target"]
        trigger = step["trigger"]
        if anim not in SUPPORTED_ANIMATIONS:
            results.append(StepResult(
                target=target, animation=anim, injected_xml=False,
                note="animation type not supported by injector — speaker notes only",
            ))
            continue
        if trigger not in SUPPORTED_TRIGGERS:
            results.append(StepResult(
                target=target, animation=anim, injected_xml=False,
                note=f"trigger '{trigger}' too fragile to author by hand — speaker notes only",
            ))
            continue
        if target not in shape_id_lookup:
            results.append(StepResult(
                target=target, animation=anim, injected_xml=False,
                note=f"shape name '{target}' not found on slide — speaker notes only",
            ))
            continue
        injectable_steps.append(step)
        results.append(StepResult(
            target=target, animation=anim, injected_xml=True,
            note="injected",
        ))

    if not injectable_steps:
        return results

    timing = build_timing_xml(injectable_steps, shape_id_lookup)

    # Remove any existing timing element, then append new one to <p:sld>
    slide_element = slide._element
    existing_timing = slide_element.find(qn("p:timing"))
    if existing_timing is not None:
        slide_element.remove(existing_timing)
    slide_element.append(timing)

    return results


def main() -> int:
    description = (__doc__ or "Inject PowerPoint animations.").split("\n\n")[0]
    parser = argparse.ArgumentParser(description=description)
    parser.add_argument("--in", dest="input_path", required=True, type=Path)
    parser.add_argument("--spec", required=True, type=Path)
    parser.add_argument("--out", required=True, type=Path)
    args = parser.parse_args()

    spec = json.loads(args.spec.read_text())
    prs = Presentation(str(args.input_path))

    total_steps = 0
    total_injected = 0
    print(f"Loaded {args.input_path}")
    print(f"  Slides: {len(prs.slides)}")
    print(f"  Spec covers: {len(spec.get('slides', []))} slide(s)")
    print()

    for slide_spec in spec.get("slides", []):
        idx = slide_spec["slide_index"]
        steps = slide_spec.get("steps", [])
        if idx >= len(prs.slides):
            print(f"WARN: slide_index {idx} out of range, skipping")
            continue
        slide = prs.slides[idx]
        results = inject_for_slide(slide, steps)
        print(f"Slide {idx + 1}: {len(steps)} steps")
        for r in results:
            marker = "[ok]   " if r.injected_xml else "[notes]"
            print(f"  {marker} {r.target} ({r.animation}) -- {r.note}")
            total_steps += 1
            if r.injected_xml:
                total_injected += 1
        print()

    prs.save(str(args.out))
    print(f"Wrote {args.out}")
    print(f"  XML-injected: {total_injected}/{total_steps} steps")
    print(f"  Speaker-note mirror: 100% (always written)")
    print()
    print("Verify by opening the .pptx in PowerPoint or LibreOffice Impress.")
    print("If XML-injected animations don't appear or look wrong, apply them")
    print("manually using the speaker notes as the source of truth.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
